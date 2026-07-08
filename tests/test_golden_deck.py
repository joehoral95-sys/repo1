"""End-to-end golden deck build: the CI gate for engine regressions."""

from pptx import Presentation

from deckstudio.validate.build_check import check_build


def test_build_succeeds(built_example, example_spec):
    assert built_example.output_path.exists()
    assert built_example.slide_count == len(example_spec.slides)
    assert built_example.warnings == [], built_example.warnings


def test_smoke_check_passes(built_example, example_spec):
    problems = check_build(built_example.output_path, example_spec)
    assert problems == []


def test_every_chart_is_natively_editable(built_example):
    """THE guarantee: every chart has an embedded workbook (Edit Data works)."""
    prs = Presentation(str(built_example.output_path))
    charts = [sh for sl in prs.slides for sh in sl.shapes if getattr(sh, "has_chart", False)]
    assert charts, "golden deck should contain charts"
    for shape in charts:
        xlsx_part = shape.chart.part.chart_workbook.xlsx_part
        assert xlsx_part is not None
        assert len(xlsx_part.blob) > 0


def test_chart_values_round_trip(built_example, example_spec):
    """Values written to the chart can be read back — proves real data binding."""
    prs = Presentation(str(built_example.output_path))
    chart_slides = {i: s for i, s in enumerate(example_spec.slides) if s.type == "chart"}
    for i, model in chart_slides.items():
        shapes = [sh for sh in prs.slides[i].shapes if getattr(sh, "has_chart", False)]
        assert len(shapes) == 1
        chart = shapes[0].chart
        plot = chart.plots[0]
        got = [list(s.values) for s in plot.series]
        want = [s.values for s in model.chart.series]
        assert got == want, f"slide '{model.id}' chart values changed in round-trip"


def test_versioning_never_overwrites(example_spec, tmp_path):
    from deckstudio.engine.builder import build_deck

    r1 = build_deck(example_spec, output_dir=tmp_path, deck_name="d")
    r2 = build_deck(example_spec, output_dir=tmp_path, deck_name="d")
    assert r1.version == 1 and r2.version == 2
    assert r1.output_path.exists() and r2.output_path.exists()
    assert r1.output_path != r2.output_path


def test_footers_and_page_numbers(built_example, example_spec):
    prs = Presentation(str(built_example.output_path))
    # slide 4 (index 3) is a big_number content slide -> has footer + page number
    texts = [sh.text_frame.text for sh in prs.slides[3].shapes if sh.has_text_frame]
    assert any("Society of Actuaries" in t for t in texts)
    assert any(t.strip() == "4" for t in texts)


def test_notes_written(built_example, example_spec):
    prs = Presentation(str(built_example.output_path))
    for i, model in enumerate(example_spec.slides):
        if model.notes:
            assert prs.slides[i].has_notes_slide
            note = prs.slides[i].notes_slide.notes_text_frame.text
            assert model.notes.strip()[:30] in note
