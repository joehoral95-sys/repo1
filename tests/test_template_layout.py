"""The engine must never build on a decorated layout.

The SOA template ships 'blank-ish' layouts (e.g. '1_Blank ') carrying a
full-slide WMF pattern that PowerPoint renders BEHIND slide content — and
LibreOffice previews don't show WMF at all, so this only surfaces on a real
machine. These tests lock in the fix.
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckstudio.engine.template import blank_layout, open_template
from deckstudio.tokens import load_tokens


def _decorations(layout):
    return [sh for sh in layout.shapes
            if sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)]


def test_brand_template_blank_layout_is_undecorated():
    prs = open_template(load_tokens().brand_dir / "template.pptx")
    layout = blank_layout(prs)
    assert layout.name.strip().lower() == "blank"
    assert _decorations(layout) == [], (
        f"layout '{layout.name}' carries background art that would render "
        "behind slide content in PowerPoint")


def test_built_slides_use_undecorated_layout(built_example):
    from pptx import Presentation

    prs = Presentation(str(built_example.output_path))
    for i, slide in enumerate(prs.slides, start=1):
        assert _decorations(slide.slide_layout) == [], (
            f"slide {i} sits on decorated layout '{slide.slide_layout.name}'")
