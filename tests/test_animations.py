"""Animations are on by default (subtle) and disabled per deck with
animations: off — required for decks read on phones/web previews, whose
renderers mis-composite animated slides."""

from lxml import etree
from pptx import Presentation

from deckstudio.engine.builder import build_deck
from deckstudio.spec.schema import DeckSpec

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

ANIMATED_SPEC = {
    "deck": {"title": "T", "audience": "a", "intent": "i", "animations": "subtle"},
    "slides": [
        {"id": "cover", "type": "title", "title": "Hello"},
        {"id": "stats", "type": "big_number", "title": "Numbers went up this quarter",
         "animate": "build",
         "stats": [{"value": "1", "label": "one"}, {"value": "2", "label": "two"}]},
        {"id": "plain", "type": "quote", "text": "No animation here."},
        {"id": "soft", "type": "content", "title": "Context arrived all at once",
         "animate": "fade", "bullets": ["a point", "another point"]},
    ],
}


def _timing_elements(pptx_path):
    prs = Presentation(str(pptx_path))
    return [sl._element.find(f"{{{P_NS}}}timing") for sl in prs.slides]


def test_golden_deck_animates_marked_slides_only(built_example, example_spec):
    timings = _timing_elements(built_example.output_path)
    for i, model in enumerate(example_spec.slides):
        if model.animate != "none":
            assert timings[i] is not None, f"slide '{model.id}' should have timing"
        else:
            assert timings[i] is None, f"slide '{model.id}' should NOT have timing"


def test_opt_in_animations_present_only_on_animated_slides(tmp_path):
    spec = DeckSpec.model_validate(ANIMATED_SPEC)
    result = build_deck(spec, output_dir=tmp_path, deck_name="anim")
    timings = _timing_elements(result.output_path)
    assert timings[1] is not None, "build slide should carry timing"
    assert timings[3] is not None, "fade slide should carry timing"
    assert timings[0] is None and timings[2] is None


def test_timing_xml_is_wellformed_and_targets_real_shapes(tmp_path):
    spec = DeckSpec.model_validate(ANIMATED_SPEC)
    result = build_deck(spec, output_dir=tmp_path, deck_name="anim2")
    prs = Presentation(str(result.output_path))
    found = 0
    for slide in prs.slides:
        timing = slide._element.find(f"{{{P_NS}}}timing")
        if timing is None:
            continue
        found += 1
        etree.fromstring(etree.tostring(timing))  # well-formed
        shape_ids = {str(sh.shape_id) for sh in slide.shapes}
        spids = {el.get("spid") for el in timing.iter(f"{{{P_NS}}}spTgt")}
        assert spids and spids <= shape_ids
    assert found == 2


def test_cli_kill_switch(tmp_path):
    spec = DeckSpec.model_validate(ANIMATED_SPEC)
    result = build_deck(spec, output_dir=tmp_path, deck_name="killed",
                        enable_animations=False)
    assert all(t is None for t in _timing_elements(result.output_path))
