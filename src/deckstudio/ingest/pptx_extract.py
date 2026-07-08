"""The 'beautify' entry point: deterministically extract content from an
existing deck into a draft spec + gap report.

Extraction is mechanical (every text slide becomes `type: content`; native
charts and tables port losslessly). The AGENT then does the judgment step:
re-outline, pick better slide types, add insight framing, and produce the
real spec.yaml. See AGENTS.md step "Beautify an existing deck".
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ruamel.yaml import YAML


def extract_pptx(path: Path, deck_dir: Path) -> tuple[Path, Path]:
    """Write extracted.yaml (draft spec) and extracted_report.md into deck_dir.

    Returns (yaml_path, report_path)."""
    prs = Presentation(str(path))
    slides_out: list[dict] = []
    report: list[str] = [
        f"# Extraction report: {path.name}",
        "",
        f"{len(prs.slides)} slides scanned. Items below need a human/agent decision.",
        "",
    ]
    media_dir = deck_dir / "sources" / "extracted_media"

    for i, slide in enumerate(prs.slides, start=1):
        entry: dict = {"id": f"slide-{i}", "type": "content"}
        bullets: list[str] = []
        title = None
        flags: list[str] = []

        for shape in slide.shapes:
            title, bullets, extra = _walk_shape(
                shape, slide, title, bullets, entry, media_dir, i)
            flags.extend(extra)

        entry["title"] = title or f"(untitled slide {i})"
        if "chart" in entry:
            entry["type"] = "chart"
            entry.pop("bullets", None)
        elif "table_data" in entry:
            entry["type"] = "table"
            table = entry.pop("table_data")
            entry["columns"] = table[0]
            entry["rows"] = table[1:] if len(table) > 1 else []
        elif bullets:
            entry["bullets"] = bullets
        else:
            entry["bullets"] = ["(no text found on this slide)"]
            flags.append("no text content found")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                entry["notes"] = notes

        if flags:
            report.append(f"- **Slide {i}** ({entry['title'][:60]}): " + "; ".join(flags))
        slides_out.append(entry)

    draft = {
        "deck": {
            "title": path.stem,
            "audience": "TODO — ask Joe",
            "intent": "TODO — ask Joe",
        },
        "slides": slides_out,
    }

    deck_dir.mkdir(parents=True, exist_ok=True)
    yaml_path = deck_dir / "extracted.yaml"
    yaml = YAML()
    yaml.default_flow_style = False
    yaml.width = 100
    with yaml_path.open("w", encoding="utf-8") as f:
        yaml.dump(draft, f)

    if len(report) == 4:
        report.append("- Nothing flagged: all content extracted cleanly.")
    report_path = deck_dir / "extracted_report.md"
    report_path.write_text("\n".join(report) + "\n", encoding="utf-8")
    return yaml_path, report_path


def _walk_shape(shape, slide, title, bullets, entry, media_dir, slide_no):
    """Extract one shape; returns (title, bullets, flags)."""
    flags: list[str] = []
    st = shape.shape_type

    if getattr(shape, "has_chart", False):
        chart = shape.chart
        try:
            entry["chart"] = _chart_to_spec(chart)
        except Exception as e:
            flags.append(f"chart could not be read ({e}) — redesign it from source data")
        return title, bullets, flags

    if st == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        entry["table_data"] = [
            [cell.text.strip() for cell in row.cells] for row in table.rows]
        return title, bullets, flags

    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            media_dir.mkdir(parents=True, exist_ok=True)
            fname = f"slide{slide_no}_{shape.shape_id}.{image.ext}"
            (media_dir / fname).write_bytes(image.blob)
            flags.append(f"picture saved to sources/extracted_media/{fname} — "
                         "decide whether it earns its place")
        except Exception:
            flags.append("picture could not be exported")
        return title, bullets, flags

    if st == MSO_SHAPE_TYPE.GROUP:
        flags.append("grouped shapes (possibly SmartArt/diagram) — content below may "
                     "be partial; consider redesigning with icon_row/timeline/comparison")
        for sub in shape.shapes:
            title, bullets, extra = _walk_shape(
                sub, slide, title, bullets, entry, media_dir, slide_no)
            flags.extend(extra)
        return title, bullets, flags

    if shape.has_text_frame:
        is_title = False
        try:
            ph = shape.placeholder_format
            if ph is not None and ph.idx == 0:
                is_title = True
        except ValueError:
            pass
        text_lines = [p.strip() for p in shape.text_frame.text.split("\n") if p.strip()]
        if not text_lines:
            return title, bullets, flags
        if is_title or (title is None and len(text_lines) == 1 and len(text_lines[0]) < 100):
            title = title or text_lines[0]
            text_lines = text_lines[1:] if not is_title else text_lines[1:]
        bullets.extend(_clean(t) for t in text_lines)
    return title, bullets, flags


def _chart_to_spec(chart) -> dict:
    from pptx.enum.chart import XL_CHART_TYPE as XL

    kind_map = {
        XL.COLUMN_CLUSTERED: "bar_clustered",
        XL.COLUMN_STACKED: "bar_stacked",
        XL.BAR_CLUSTERED: "bar_horizontal",
        XL.LINE: "line",
        XL.LINE_MARKERS: "line",
        XL.AREA: "area",
        XL.PIE: "pie",
        XL.DOUGHNUT: "doughnut",
        XL.XY_SCATTER: "scatter",
    }
    kind = kind_map.get(chart.chart_type, "bar_clustered")
    plot = chart.plots[0]
    categories = [str(c) for c in plot.categories]
    series = [
        # str() everywhere: python-pptx returns str subclasses that YAML
        # serializers refuse to represent.
        {"name": str(s.name) if s.name else f"Series {j + 1}",
         "values": [float(v) if v is not None else None for v in s.values]}
        for j, s in enumerate(plot.series)
    ]
    return {"kind": kind, "categories": categories, "series": series}


def _clean(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()
