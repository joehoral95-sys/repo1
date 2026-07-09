"""Comparison: two panels; the emphasized option gets the primary treatment."""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from ...spec.schema import ComparisonSide, ComparisonSlide
from ..geometry import Box, columns, inset
from ..registry import renderer
from ..shapes import add_chip, add_rect
from ..text import add_bullets, add_text
from ._common import add_title_band


@renderer("comparison")
def render(slide, model: ComparisonSlide, ctx) -> None:
    """Variants: panels (filled cards) | open (airy, divider line only)."""
    tokens = ctx.tokens
    variant, chrome = ctx.variant_parts(model)
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker,
                          style=chrome)
    left_col, right_col = columns(area, 2, tokens.gutter_in + 0.35)
    if variant == "open":
        _open_side(slide, left_col, model.left, tokens,
                   emphasized=model.emphasize == "left")
        _open_side(slide, right_col, model.right, tokens,
                   emphasized=model.emphasize == "right")
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches((left_col.right_in + right_col.left_in) / 2 - 0.008),
            Inches(area.top_in + 0.1), Inches(0.016), Inches(area.height_in - 0.5))
        divider.fill.solid()
        divider.fill.fore_color.rgb = tokens.color("border")
        divider.line.fill.background()
        divider.shadow.inherit = False
    else:
        _panel(slide, left_col, model.left, tokens,
               emphasized=model.emphasize == "left")
        _panel(slide, right_col, model.right, tokens,
               emphasized=model.emphasize == "right")

    if variant == "open":
        if model.emphasize != "none" and model.badge:
            col = left_col if model.emphasize == "left" else right_col
            add_chip(slide, col.left_in, area.top_in - 0.05, model.badge.upper(),
                     tokens, fill="accent_warm", text_color="primary")
        return

    # "vs" disc on the seam between the panels
    d = 0.62
    cx = (left_col.right_in + right_col.left_in) / 2 - d / 2
    cy = area.top_in + area.height_in * 0.42
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy),
                                  Inches(d), Inches(d))
    disc.fill.solid()
    disc.fill.fore_color.rgb = tokens.color("accent_warm")
    disc.line.color.rgb = tokens.color("white")
    from pptx.util import Pt
    disc.line.width = Pt(2.5)
    disc.shadow.inherit = False
    add_text(slide, Box(cx, cy + 0.06, d, d - 0.12), "vs", tokens, scale="subtitle",
             role="heading", color="primary", bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # badge the emphasized panel, only when the spec names the badge
    if model.emphasize != "none" and model.badge:
        col = left_col if model.emphasize == "left" else right_col
        add_chip(slide, col.left_in + 0.35, col.top_in - 0.15, model.badge.upper(),
                 tokens, fill="accent_warm", text_color="primary")


def _panel(slide, box: Box, side: ComparisonSide, tokens, *, emphasized: bool) -> None:
    fill = "primary" if emphasized else "neutral_light"
    heading_color = "accent_warm" if emphasized else "primary"
    body_color = "white" if emphasized else "neutral_dark"
    add_rect(slide, box, tokens, fill=fill, rounded=True, corner=0.05)
    pad = inset(box, x_in=0.35, y_in=0.3)
    add_text(slide, Box(pad.left_in, pad.top_in, pad.width_in, 0.5), side.heading,
             tokens, scale="subtitle", role="heading", color=heading_color,
             bold=True, shrink_to_fit=True)
    bullets_box = Box(pad.left_in, pad.top_in + 0.65, pad.width_in,
                      pad.height_in - 0.65)
    add_bullets(slide, bullets_box, list(side.bullets), tokens, color=body_color,
                space_after_pt=8)


def _open_side(slide, box, side, tokens, *, emphasized: bool) -> None:
    """No fills: colored heading + rule, bullets on white."""
    heading_color = "accent" if emphasized else "primary"
    add_text(slide, Box(box.left_in, box.top_in + 0.35, box.width_in, 0.5),
             side.heading, tokens, scale="subtitle", role="heading",
             color=heading_color, bold=True, shrink_to_fit=True)
    from ..shapes import add_accent_bar
    add_accent_bar(slide, box.left_in, box.top_in + 0.95, 0.7, tokens,
                   color=heading_color, height_in=0.05)
    add_bullets(slide, Box(box.left_in, box.top_in + 1.25, box.width_in,
                           box.height_in - 1.4),
                list(side.bullets), tokens, space_after_pt=8)
