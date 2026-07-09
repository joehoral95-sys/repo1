"""Icon row - compositions:

  columns  icon discs above centered text columns
  rows     stacked rows: icon left, heading + text right
  cards    each item in a tinted card with its icon

All use the standard title chrome (left/centered/banner/tint).
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from ...spec.schema import IconItem, IconRowSlide
from ..geometry import Box, columns, inset
from ..geometry import rows as grid_rows
from ..registry import renderer
from ..shapes import add_picture_fitted, add_rect
from ..text import adaptive_pt, add_text
from ._common import add_title_band

DISC_STYLES = [("accent", "white"), ("primary", "white"), ("accent_warm", "primary")]


def _text_pt(tokens, model: IconRowSlide) -> int:
    """Modest bump for the body text when every item is brief."""
    return adaptive_pt(int(tokens.pt("stat_label").pt),
                       [i.text for i in model.items], [(60, 16)])


@renderer("icon_row")
def render(slide, model: IconRowSlide, ctx) -> None:
    tokens = ctx.tokens
    comp, chrome = ctx.variant_parts(model)
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker,
                          style=chrome)
    if comp == "rows":
        _rows(slide, model, ctx, area)
    elif comp == "cards":
        _cards(slide, model, ctx, area)
    else:
        _columns(slide, model, ctx, area)


def _icon_path(tokens, item: IconItem):
    if not item.icon:
        return None
    icons_dir = tokens.brand_dir / "assets" / "icons"
    for candidate in (icons_dir / item.icon, icons_dir / f"{item.icon}.png"):
        if candidate.exists():
            return candidate
    return None


def _draw_icon(slide, tokens, ctx, item: IconItem, box: Box, index: int) -> None:
    path = _icon_path(tokens, item)
    if path is None and item.icon:
        ctx.warn(f"icon '{item.icon}' not found in brand/assets/icons - using fallback disc")
    if path:
        add_picture_fitted(slide, str(path), box)
        return
    fill, initial = DISC_STYLES[index % len(DISC_STYLES)]
    disc = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(box.left_in), Inches(box.top_in),
        Inches(box.width_in), Inches(box.height_in))
    disc.fill.solid()
    disc.fill.fore_color.rgb = tokens.color(fill)
    disc.line.fill.background()
    disc.shadow.inherit = False
    add_text(slide, Box(box.left_in, box.top_in + box.height_in * 0.18,
                        box.width_in, box.height_in * 0.6),
             item.heading[:1].upper(), tokens, scale="subtitle", role="heading",
             color=initial, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def _columns(slide, model: IconRowSlide, ctx, area: Box) -> None:
    tokens = ctx.tokens
    text_pt = _text_pt(tokens, model)
    icon_d = 1.0
    for i, (item, col) in enumerate(
            zip(model.items, columns(area, len(model.items), 0.4), strict=True)):
        cx = col.left_in + (col.width_in - icon_d) / 2
        icon_top = col.top_in + 0.25
        _draw_icon(slide, tokens, ctx, item, Box(cx, icon_top, icon_d, icon_d), i)
        add_text(slide, Box(col.left_in, icon_top + icon_d + 0.2, col.width_in, 0.55),
                 item.heading, tokens, scale="subtitle", role="heading",
                 color="primary", bold=True, align=PP_ALIGN.CENTER,
                 shrink_to_fit=True)
        add_text(slide, Box(col.left_in + 0.1, icon_top + icon_d + 0.8,
                            col.width_in - 0.2, col.height_in - icon_d - 1.1),
                 item.text, tokens, scale="stat_label", color="neutral_dark",
                 align=PP_ALIGN.CENTER, line_spacing=1.15, shrink_to_fit=True,
                 size_pt=text_pt)


def _rows(slide, model: IconRowSlide, ctx, area: Box) -> None:
    tokens = ctx.tokens
    text_pt = _text_pt(tokens, model)
    icon_d = 0.85
    for i, (item, row) in enumerate(
            zip(model.items, grid_rows(area, max(len(model.items), 2), 0.2),
                strict=False)):
        icon_top = row.top_in + (row.height_in - icon_d) / 2
        _draw_icon(slide, tokens, ctx, item,
                   Box(row.left_in, icon_top, icon_d, icon_d), i)
        text_left = row.left_in + icon_d + 0.45
        add_text(slide, Box(text_left, row.top_in + row.height_in * 0.14,
                            row.width_in - icon_d - 0.45, 0.45),
                 item.heading, tokens, scale="subtitle", role="heading",
                 color="primary", bold=True, shrink_to_fit=True)
        add_text(slide, Box(text_left, row.top_in + row.height_in * 0.14 + 0.5,
                            row.width_in - icon_d - 0.45,
                            row.height_in * 0.72 - 0.5),
                 item.text, tokens, scale="stat_label", color="neutral_dark",
                 line_spacing=1.1, shrink_to_fit=True, size_pt=text_pt)


def _cards(slide, model: IconRowSlide, ctx, area: Box) -> None:
    tokens = ctx.tokens
    text_pt = _text_pt(tokens, model)
    n = len(model.items)
    icon_d = 0.8
    if n == 4:
        grid = [c for r in grid_rows(area, 2, 0.25) for c in columns(r, 2, 0.25)]
    else:
        grid = columns(area, n, 0.3)
    for i, (item, card) in enumerate(zip(model.items, grid, strict=False)):
        add_rect(slide, card, tokens, fill="neutral_light", rounded=True,
                 corner=0.07, outline="border")
        pad = inset(card, x_in=0.3, y_in=0.25)
        _draw_icon(slide, tokens, ctx, item,
                   Box(pad.left_in, pad.top_in, icon_d, icon_d), i)
        add_text(slide, Box(pad.left_in, pad.top_in + icon_d + 0.15, pad.width_in,
                            0.45),
                 item.heading, tokens, scale="stat_label", role="heading",
                 color="primary", bold=True, shrink_to_fit=True)
        add_text(slide, Box(pad.left_in, pad.top_in + icon_d + 0.6, pad.width_in,
                            pad.height_in - icon_d - 0.6),
                 item.text, tokens, scale="stat_label", color="neutral_dark",
                 line_spacing=1.1, shrink_to_fit=True, size_pt=text_pt)
