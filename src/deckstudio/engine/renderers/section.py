"""Section divider — variants:

  chip     navy: sky number chip, swash title, optional metric preview
  giant    navy: oversized sky numeral, swash title
  band     navy band left with number+title; white right with the preview
  minimal  white: huge tinted numeral, blue title
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ...spec.schema import SectionSlide
from ..geometry import SLIDE_H_IN, SLIDE_W_IN, Box
from ..registry import renderer
from ..shapes import add_accent_bar, add_brand_art, add_logo, add_rect, fill_background
from ..text import add_text

PREVIEW_STATUS = {"on_track": "positive", "watch": "warning", "pending": "neutral_mid"}


@renderer("section")
def render(slide, model: SectionSlide, ctx) -> None:
    tokens = ctx.tokens
    variant = ctx.variant(model)
    if variant == "band":
        _band(slide, model, ctx)
        return
    if variant == "minimal":
        _minimal(slide, model, ctx)
        return

    m = tokens.margin_in + 0.4
    fill_background(slide, tokens, "primary")
    add_brand_art(slide, tokens, "shield_watermark", Box(9.6, 3.2, 5.2, 5.2))

    if model.number is not None:
        if variant == "giant":
            shape = add_text(slide, Box(m, 1.0, 3.4, 2.3), f"{model.number:02d}",
                             tokens, scale="big_number", role="heading",
                             color="accent_warm", bold=True)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(110)
        else:  # chip
            chip = Box(m, 1.15, 1.7, 1.7)
            add_rect(slide, chip, tokens, fill="accent_warm", rounded=True, corner=0.18)
            shape = add_text(slide, Box(chip.left_in, chip.top_in + 0.28,
                                        chip.width_in, 1.2),
                             f"{model.number:02d}", tokens, scale="big_number",
                             role="heading", color="primary", bold=True,
                             align=PP_ALIGN.CENTER)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(64)

    text_w = (6.4 if model.preview else SLIDE_W_IN - 2 * m)
    kicker = model.kicker or (f"SECTION {model.number:02d}" if model.number else None)
    if kicker:
        add_text(slide, Box(m, 3.32, text_w, 0.32), kicker.upper(), tokens,
                 scale="kicker", color="accent_warm", bold=True,
                 letter_spacing_pt=1.6)
    add_text(slide, Box(m, 3.7, text_w, 1.1), model.title, tokens,
             scale="section_title", role="heading", color="white", bold=True,
             anchor=MSO_ANCHOR.TOP, shrink_to_fit=True)
    if add_brand_art(slide, tokens, "swash_sky", Box(m, 4.75, 3.0, 0.41)) is None:
        add_accent_bar(slide, m, 4.85, 1.1, tokens, color="accent_warm", height_in=0.07)
    if model.subtitle:
        add_text(slide, Box(m, 5.35, text_w, 0.9), model.subtitle, tokens,
                 scale="subtitle", color="white", line_spacing=1.15,
                 shrink_to_fit=True)
    if model.preview:
        _preview_panel(slide, model.preview, tokens, Box(7.6, 1.5, 4.9, 4.6),
                       dark=True)
    add_logo(slide, tokens, dark_bg=True, left_in=m, bottom_in=7.0, height_in=0.4)


def _band(slide, model: SectionSlide, ctx) -> None:
    tokens = ctx.tokens
    m = tokens.margin_in
    band_w = 5.4
    add_rect(slide, Box(0, 0, band_w, SLIDE_H_IN), tokens, fill="primary",
             rounded=False)
    if model.number is not None:
        chip = Box(m + 0.2, 1.0, 1.35, 1.35)
        add_rect(slide, chip, tokens, fill="accent_warm", rounded=True, corner=0.2)
        add_text(slide, Box(chip.left_in, chip.top_in + 0.22, chip.width_in, 0.95),
                 f"{model.number:02d}", tokens, scale="section_title",
                 role="heading", color="primary", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Box(m + 0.2, 2.85, band_w - m - 0.7, 1.6), model.title, tokens,
             scale="section_title", role="heading", color="white", bold=True,
             shrink_to_fit=True)
    add_brand_art(slide, tokens, "swash_sky", Box(m + 0.2, 4.35, 2.4, 0.33))
    if model.subtitle:
        add_text(slide, Box(m + 0.2, 4.95, band_w - m - 0.7, 1.3), model.subtitle,
                 tokens, scale="subtitle", color="white", line_spacing=1.15,
                 shrink_to_fit=True)
    add_logo(slide, tokens, dark_bg=True, center_at_in=band_w / 2,
             bottom_in=7.0, height_in=0.4)
    if model.preview:
        _preview_panel(slide, model.preview, tokens,
                       Box(band_w + 0.7, 1.5, SLIDE_W_IN - band_w - 0.7 - m, 4.6),
                       dark=False)


def _minimal(slide, model: SectionSlide, ctx) -> None:
    tokens = ctx.tokens
    m = tokens.margin_in + 0.4
    if model.number is not None:
        shape = add_text(slide, Box(SLIDE_W_IN - 5.4, 1.5, 4.6, 4.2),
                         f"{model.number:02d}", tokens, scale="big_number",
                         role="heading", color="accent_tint", bold=True,
                         align=PP_ALIGN.RIGHT)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(260)
    kicker = model.kicker or (f"SECTION {model.number:02d}" if model.number else None)
    if kicker:
        add_text(slide, Box(m, 3.0, 8.0, 0.32), kicker.upper(), tokens,
                 scale="kicker", color="primary", bold=True, letter_spacing_pt=1.6)
    add_text(slide, Box(m, 3.38, 8.6, 1.1), model.title, tokens,
             scale="section_title", role="heading", color="accent", bold=True,
             anchor=MSO_ANCHOR.TOP, shrink_to_fit=True)
    add_brand_art(slide, tokens, "swash_sky", Box(m, 4.45, 3.0, 0.41))
    if model.subtitle:
        add_text(slide, Box(m, 5.05, 8.6, 0.9), model.subtitle, tokens,
                 scale="subtitle", color="neutral_dark", line_spacing=1.15,
                 shrink_to_fit=True)
    add_logo(slide, tokens, left_in=m, bottom_in=7.05, height_in=0.38)


def _preview_panel(slide, metrics, tokens, panel: Box, *, dark: bool) -> None:
    """Metric teaser rows with status dots — on navy (dark) or white."""
    from pptx.enum.shapes import MSO_SHAPE

    label_color = "white" if dark else "neutral_dark"
    value_color = "accent_warm" if dark else "accent"
    rule_color = RGBColor.from_string("2A4A6B") if dark else tokens.color("border")
    n = len(metrics)
    row_h = min(0.75, panel.height_in / n)
    for i, metric in enumerate(metrics):
        y = panel.top_in + i * row_h
        if i > 0:
            rule = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(panel.left_in), Inches(y),
                Inches(panel.width_in), Inches(0.011))
            rule.fill.solid()
            rule.fill.fore_color.rgb = rule_color
            rule.line.fill.background()
            rule.shadow.inherit = False
        if metric.status:
            d = 0.11
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(panel.left_in),
                Inches(y + (row_h - d) / 2), Inches(d), Inches(d))
            dot.fill.solid()
            dot.fill.fore_color.rgb = tokens.color(PREVIEW_STATUS[metric.status])
            dot.line.fill.background()
            dot.shadow.inherit = False
        add_text(slide, Box(panel.left_in + 0.3, y, panel.width_in - 1.7, row_h),
                 metric.label, tokens, scale="stat_label", color=label_color,
                 anchor=MSO_ANCHOR.MIDDLE, shrink_to_fit=True)
        add_text(slide, Box(panel.right_in - 1.35, y, 1.35, row_h), metric.value,
                 tokens, scale="stat_label", color=value_color, bold=True,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
