"""Table: native PowerPoint table with brand styling (no default-blue theme)."""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ...spec.schema import TableSlide
from ..registry import renderer
from ..text import set_runs
from ._common import add_title_band


@renderer("table")
def render(slide, model: TableSlide, ctx) -> None:
    tokens = ctx.tokens
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker)
    nrows = len(model.rows) + 1
    ncols = len(model.columns)
    # Short tables get roomier rows so they don't huddle under the title.
    row_h = 0.58 if len(model.rows) <= 5 else 0.42
    height = min(area.height_in, 0.55 + row_h * len(model.rows))
    shape = slide.shapes.add_table(nrows, ncols, area.left, area.top,
                                   area.width, Inches(height))
    table = shape.table
    table.first_row = False
    table.horz_banding = False

    body_font = tokens.body_font.name
    for c, col_name in enumerate(model.columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = tokens.color("primary")
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        set_runs(p, col_name, font=body_font, size=tokens.pt("table_header"),
                 color=tokens.color("white"), accent_color=tokens.color("accent_warm"),
                 bold=True)

    for r, row in enumerate(model.rows, start=1):
        banded = r % 2 == 0
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            emphasized = model.emphasize_col is not None and c == model.emphasize_col
            if emphasized:
                # soft highlight: pale sky fill + bold blue text, so the
                # column draws the eye without shouting
                cell.fill.fore_color.rgb = tokens.color("accent_tint")
            else:
                cell.fill.fore_color.rgb = (tokens.color("neutral_light") if banded
                                            else tokens.color("white"))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            set_runs(p, value, font=body_font, size=tokens.pt("table_body"),
                     color=tokens.color("accent") if emphasized else tokens.color("neutral_dark"),
                     accent_color=tokens.color("accent"),
                     bold=c == 0 or emphasized)
