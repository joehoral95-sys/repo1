"""Chart slide: assertion title + native chart + the takeaway.

When the spec provides an `insight`, the slide uses the exhibit layout:
chart on the left two-thirds, a navy takeaway panel on the right carrying
the "so what" — visually distinct from every other slide family. Without an
insight the chart takes the full width.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from ...spec.schema import ChartSlide
from ..charts import add_chart
from ..geometry import Box, inset
from ..registry import renderer
from ..shapes import add_accent_bar, add_chip, add_rect, chip_width
from ..text import add_text
from ._common import add_title_band


@renderer("chart")
def render(slide, model: ChartSlide, ctx) -> None:
    tokens = ctx.tokens
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker)

    if model.insight:
        panel_w = 3.4
        chart_box = Box(area.left_in, area.top_in,
                        area.width_in - panel_w - 0.35, area.height_in - 0.3)
        add_chart(slide, model.chart, chart_box, tokens)
        _annotate_highlight(slide, model, chart_box, tokens)
        _takeaway_panel(slide, model, tokens,
                        Box(area.right_in - panel_w, area.top_in, panel_w, area.height_in - 0.3))
        return

    bottom_pad = 0.3 if model.source else 0.0
    add_chart(slide, model.chart, Box(area.left_in, area.top_in, area.width_in,
                                      area.height_in - bottom_pad), tokens)
    if model.source:
        add_text(slide, Box(area.left_in, area.bottom_in - 0.22, area.width_in, 0.22),
                 f"Source: {model.source}", tokens, scale="caption",
                 color="neutral_mid", align=PP_ALIGN.RIGHT)


def _annotate_highlight(slide, model: ChartSlide, chart_box: Box, tokens) -> None:
    """When a point is highlighted, pin an auto-computed delta chip to the
    chart's top-right — the editorial 'look here' annotation."""
    hl = model.chart.highlight
    if hl is None or hl.point == 0:
        return
    series = next(s for s in model.chart.series if s.name == hl.series)
    curr, prev = series.values[hl.point], series.values[hl.point - 1]
    if curr is None or prev is None or prev == 0:
        return
    pct = (curr - prev) / abs(prev) * 100
    arrow = "▼" if pct < 0 else "▲"
    label = model.chart.categories[hl.point]
    text = f"{label}: {arrow} {abs(pct):.1f}%"
    width_est = chip_width(text)
    add_chip(slide, chart_box.right_in - width_est - 0.1, chart_box.top_in + 0.02,
             text, tokens, fill="highlight", text_color="primary")


def _takeaway_panel(slide, model: ChartSlide, tokens, box: Box) -> None:
    """Navy card: SO WHAT kicker, the insight, source at the bottom."""
    add_rect(slide, box, tokens, fill="primary", rounded=True, corner=0.05)
    pad = inset(box, x_in=0.3, y_in=0.35)
    add_text(slide, Box(pad.left_in, pad.top_in, pad.width_in, 0.3), "SO WHAT",
             tokens, scale="kicker", color="accent_warm", bold=True)
    add_accent_bar(slide, pad.left_in, pad.top_in + 0.38, 0.55, tokens,
                   color="accent_warm", height_in=0.05)
    add_text(slide, Box(pad.left_in, pad.top_in + 0.6, pad.width_in,
                        pad.height_in - (1.0 if model.source else 0.6)),
             model.insight, tokens, scale="subtitle", role="heading",
             color="white", bold=True, line_spacing=1.15, shrink_to_fit=True)
    if model.source:
        add_text(slide, Box(pad.left_in, pad.bottom_in - 0.4, pad.width_in, 0.4),
                 f"Source: {model.source}", tokens, scale="caption",
                 color="neutral_light", shrink_to_fit=True)
