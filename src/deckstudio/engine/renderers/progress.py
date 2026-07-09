"""Progress rings: 'how far along are we' rendered as donut gauges."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from ...spec.schema import ProgressSlide
from ..geometry import Box, columns, vcenter
from ..infographics import add_progress_ring
from ..registry import renderer
from ..text import add_text
from ._common import add_title_band


@renderer("progress")
def render(slide, model: ProgressSlide, ctx) -> None:
    """Variants: rings (donut gauges) | bars (horizontal percent bars)."""
    tokens = ctx.tokens
    comp, chrome = ctx.variant_parts(model)
    if comp == "bars":
        _bars(slide, model, ctx, chrome)
        return
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker,
                          style=chrome)
    n = len(model.items)
    caption_pad = 0.55 if model.caption else 0.0
    ring_area = Box(area.left_in, area.top_in, area.width_in,
                    area.height_in - caption_pad)
    ring_h = min(3.1, ring_area.height_in - 0.2)
    for item, col in zip(model.items, columns(ring_area, n, 0.5), strict=True):
        add_progress_ring(slide, vcenter(col, ring_h), item.percent / 100.0,
                          item.label, tokens)
    if model.caption:
        add_text(slide, Box(area.left_in, area.bottom_in - 0.4, area.width_in, 0.4),
                 model.caption, tokens, scale="stat_label", color="neutral_mid",
                 align=PP_ALIGN.CENTER, shrink_to_fit=True)


def _bars(slide, model: ProgressSlide, ctx, chrome: str = "left") -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Inches

    from ..geometry import rows as grid_rows

    tokens = ctx.tokens
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker,
                          style=chrome)
    caption_pad = 0.55 if model.caption else 0.0
    bar_area = Box(area.left_in, area.top_in + 0.3, area.width_in,
                   area.height_in - caption_pad - 0.5)
    label_w = 3.1
    track_w = bar_area.width_in - label_w - 1.15
    for item, row in zip(model.items,
                         grid_rows(bar_area, max(len(model.items), 3), 0.25),
                         strict=False):
        bar_h = min(0.5, row.height_in * 0.5)
        bar_top = row.top_in + (row.height_in - bar_h) / 2
        add_text(slide, Box(bar_area.left_in, bar_top - 0.05, label_w - 0.2,
                            bar_h + 0.1),
                 item.label, tokens, scale="stat_label", color="neutral_dark",
                 anchor=MSO_ANCHOR.MIDDLE, shrink_to_fit=True)
        track = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_area.left_in + label_w),
            Inches(bar_top), Inches(track_w), Inches(bar_h))
        track.adjustments[0] = 0.5
        track.fill.solid()
        track.fill.fore_color.rgb = tokens.color("neutral_light")
        track.line.fill.background()
        track.shadow.inherit = False
        frac = max(0.0, min(1.0, item.percent / 100.0))
        if frac > 0:
            fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_area.left_in + label_w),
                Inches(bar_top), Inches(max(0.25, track_w * frac)), Inches(bar_h))
            fill.adjustments[0] = 0.5
            fill.fill.solid()
            fill.fill.fore_color.rgb = tokens.color("accent")
            fill.line.fill.background()
            fill.shadow.inherit = False
        add_text(slide, Box(bar_area.left_in + label_w + track_w + 0.15,
                            bar_top - 0.05, 1.0, bar_h + 0.1),
                 f"{round(item.percent)}%", tokens, scale="subtitle",
                 role="heading", color="primary", bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    if model.caption:
        add_text(slide, Box(area.left_in, area.bottom_in - 0.4, area.width_in, 0.4),
                 model.caption, tokens, scale="stat_label", color="neutral_mid",
                 align=PP_ALIGN.CENTER, shrink_to_fit=True)
