"""Quote / big-statement slide: white card floating on the navy gradient,
giant sky quote mark breaking the card's top edge — editorial, not flat."""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from ...spec.schema import QuoteSlide
from ..geometry import SLIDE_W_IN, Box, inset
from ..registry import renderer
from ..shapes import add_rect, fill_background
from ..text import add_text


@renderer("quote")
def render(slide, model: QuoteSlide, ctx) -> None:
    """Variants: card (white card on navy) | dark (bare on navy) | light."""
    tokens = ctx.tokens
    variant = ctx.variant(model)
    if variant in ("dark", "light"):
        _bare(slide, model, ctx, dark=(variant == "dark"))
        return
    fill_background(slide, tokens, "primary")

    card = Box(1.55, 1.95, SLIDE_W_IN - 3.1, 3.15)
    add_rect(slide, card, tokens, fill="white", rounded=True, corner=0.05)

    pad = inset(card, x_in=0.7, y_in=0.45)
    add_text(slide, pad, model.text, tokens, scale="quote", role="heading",
             color="primary", bold=True, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.15, shrink_to_fit=True)

    # oversized quote mark straddling the card's top-left corner
    mark = add_text(slide, Box(card.left_in - 0.35, card.top_in - 1.25, 2.0, 1.9),
                    "“", tokens, scale="big_number", role="heading",
                    color="accent_warm", bold=True)
    for para in mark.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(130)

    if model.attribution:
        add_text(slide, Box(card.left_in + 0.1, card.bottom_in + 0.35,
                            card.width_in - 0.2, 0.5),
                 f"—  {model.attribution}", tokens, scale="subtitle",
                 color="accent_warm", bold=True)


def _bare(slide, model: QuoteSlide, ctx, *, dark: bool) -> None:
    tokens = ctx.tokens
    if dark:
        fill_background(slide, tokens, "primary")
    m = tokens.margin_in + 1.0
    mark = add_text(slide, Box(m - 0.5, 0.9, 2.0, 1.6), "\u201c", tokens,
                    scale="big_number", role="heading", color="accent_warm",
                    bold=True)
    for para in mark.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(120)
    add_text(slide, Box(m, 2.35, SLIDE_W_IN - 2 * m, 2.6), model.text, tokens,
             scale="quote", role="heading",
             color="white" if dark else "primary",
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, line_spacing=1.15,
             shrink_to_fit=True)
    if model.attribution:
        add_text(slide, Box(m, 5.3, SLIDE_W_IN - 2 * m, 0.6),
                 f"\u2014  {model.attribution}", tokens, scale="subtitle",
                 color="accent_warm" if dark else "accent", bold=True)
