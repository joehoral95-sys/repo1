"""Timeline: horizontal spine with milestone nodes; done nodes are filled."""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ...spec.schema import TimelineSlide
from ..geometry import Box
from ..registry import renderer
from ..text import adaptive_pt, add_text
from ._common import add_title_band


@renderer("timeline")
def render(slide, model: TimelineSlide, ctx) -> None:
    """Variants: spine (horizontal) | vertical (stacked rows, long labels)."""
    tokens = ctx.tokens
    comp, chrome = ctx.variant_parts(model)
    if comp == "vertical":
        _vertical(slide, model, ctx, chrome)
        return
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker,
                          style=chrome)
    n = len(model.milestones)
    spine_y = area.top_in + area.height_in * 0.42
    node_d = 0.34

    spine = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(area.left_in + 0.2), Inches(spine_y - 0.015),
        Inches(area.width_in - 0.4), Inches(0.03))
    spine.fill.solid()
    spine.fill.fore_color.rgb = tokens.color("neutral_mid")
    spine.line.fill.background()
    spine.shadow.inherit = False

    step = (area.width_in - 0.4 - node_d) / max(n - 1, 1)
    next_idx = next((i for i, ms in enumerate(model.milestones) if not ms.done), None)
    # Short labels earn bigger type (one size for the whole spine).
    label_pt = adaptive_pt(int(tokens.pt("stat_label").pt),
                           [ms.label for ms in model.milestones],
                           [(12, 18), (20, 16)])
    date_pt = adaptive_pt(int(tokens.pt("caption").pt),
                          [ms.date or "" for ms in model.milestones],
                          [(10, 14)])
    for i, ms in enumerate(model.milestones):
        cx = area.left_in + 0.2 + i * step
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx), Inches(spine_y - node_d / 2),
            Inches(node_d), Inches(node_d))
        node.shadow.inherit = False
        if ms.done:
            node.fill.solid()
            node.fill.fore_color.rgb = tokens.color("accent")
            node.line.fill.background()
        elif i == next_idx:
            # "you are here": the next milestone pops in sky blue
            node.fill.solid()
            node.fill.fore_color.rgb = tokens.color("accent_warm")
            node.line.fill.background()
        else:
            node.fill.solid()
            node.fill.fore_color.rgb = tokens.color("white")
            node.line.color.rgb = tokens.color("accent")
            node.line.width = Pt(2)

        label_w = min(step + node_d, 2.4)
        label_left = cx + node_d / 2 - label_w / 2
        label_left = min(max(label_left, area.left_in), area.right_in - label_w)
        above = i % 2 == 0
        # Date hugs the spine; the label is anchored to grow AWAY from the
        # date so multi-line labels can never collide with it.
        if ms.date:
            date_top = spine_y - 0.58 if above else spine_y + 0.3
            add_text(slide, Box(label_left, date_top, label_w, 0.3), ms.date, tokens,
                     scale="caption", color="accent", bold=True, align=PP_ALIGN.CENTER,
                     size_pt=date_pt)
        if above:
            label_box = Box(label_left, spine_y - 1.85, label_w, 1.2)
            anchor = MSO_ANCHOR.BOTTOM
        else:
            label_box = Box(label_left, spine_y + 0.64, label_w, 1.2)
            anchor = MSO_ANCHOR.TOP
        add_text(slide, label_box, ms.label, tokens,
                 scale="stat_label", color="neutral_dark", align=PP_ALIGN.CENTER,
                 anchor=anchor, shrink_to_fit=True, size_pt=label_pt)


def _vertical(slide, model: TimelineSlide, ctx, chrome: str = "left") -> None:
    from ..geometry import rows as grid_rows

    tokens = ctx.tokens
    area = add_title_band(slide, tokens, model.title, kicker=model.kicker,
                          style=chrome)
    node_d = 0.3
    rail_x = area.left_in + 2.15
    next_idx = next((i for i, ms in enumerate(model.milestones) if not ms.done), None)
    label_pt = adaptive_pt(int(tokens.pt("body").pt),
                           [ms.label for ms in model.milestones],
                           [(16, 20), (30, 18)])
    mrows = grid_rows(area, max(len(model.milestones), 3), 0.12)
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(rail_x + node_d / 2 - 0.015),
        Inches(mrows[0].top_in + mrows[0].height_in / 2),
        Inches(0.03),
        Inches(mrows[len(model.milestones) - 1].top_in
               + mrows[len(model.milestones) - 1].height_in / 2
               - mrows[0].top_in - mrows[0].height_in / 2))
    rail.fill.solid()
    rail.fill.fore_color.rgb = tokens.color("neutral_mid")
    rail.line.fill.background()
    rail.shadow.inherit = False
    for i, ms in enumerate(model.milestones):
        row = mrows[i]
        cy = row.top_in + row.height_in / 2
        if ms.date:
            add_text(slide, Box(area.left_in, cy - 0.18, 1.9, 0.36), ms.date,
                     tokens, scale="stat_delta", color="accent", bold=True,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(rail_x), Inches(cy - node_d / 2),
            Inches(node_d), Inches(node_d))
        node.shadow.inherit = False
        if ms.done:
            node.fill.solid()
            node.fill.fore_color.rgb = tokens.color("accent")
            node.line.fill.background()
        elif i == next_idx:
            node.fill.solid()
            node.fill.fore_color.rgb = tokens.color("accent_warm")
            node.line.fill.background()
        else:
            node.fill.solid()
            node.fill.fore_color.rgb = tokens.color("white")
            node.line.color.rgb = tokens.color("accent")
            node.line.width = Pt(2)
        add_text(slide, Box(rail_x + node_d + 0.35, cy - 0.3,
                            area.right_in - rail_x - node_d - 0.35, 0.6),
                 ms.label, tokens, scale="body", color="neutral_dark",
                 anchor=MSO_ANCHOR.MIDDLE, shrink_to_fit=True, size_pt=label_pt)
