"""Template handling: open brand/template.pptx (or a clean 16:9 base) and
resolve slide layouts BY NAME so template edits can't silently break builds.

v1 renderers draw brand elements themselves from tokens, so they only need a
blank layout. When brand intake replaces template.pptx with SOA's official
masters, renderers can opt into named layouts (e.g. "Title Slide") and the
lookup below fails loudly if a named layout disappears.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from .geometry import SLIDE_H, SLIDE_W


class TemplateError(Exception):
    pass


BLANK_CANDIDATES = ("blank",)  # matched against name.strip().lower()


def open_template(template_path: Path | None):
    """Return a Presentation sized 16:9, from the brand template if present."""
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()  # python-pptx built-in default
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def layout_by_name(prs, name: str):
    """Find a layout by exact name across all masters; loud failure lists options."""
    available = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
            available.append(layout.name)
    raise TemplateError(
        f"Slide layout '{name}' not found in template. Available layouts: {available}. "
        "Someone may have renamed or deleted it in brand/template.pptx.")


def blank_layout(prs):
    """The layout renderers draw on. It must be genuinely blank: corporate
    templates often ship decorated 'blank-ish' layouts whose background art
    (EMF/WMF patterns) renders behind everything in PowerPoint — and preview
    tools may not show it. So we require zero picture/graphic shapes, prefer
    a layout literally named 'Blank' (whitespace/case-insensitive), and fall
    back to the least-decorated layout."""
    candidates = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            decorations = sum(
                1 for sh in layout.shapes
                if not getattr(sh, "is_placeholder", False) and _is_graphic(sh))
            candidates.append((layout, decorations, len(layout.placeholders)))
    if not candidates:
        raise TemplateError("Template has no slide layouts at all.")

    named_clean = [c for c in candidates
                   if c[0].name.strip().lower() in BLANK_CANDIDATES and c[1] == 0]
    if named_clean:
        return named_clean[0][0]

    candidates.sort(key=lambda c: (c[1], c[2]))
    layout, decorations, _ = candidates[0]
    if decorations > 0:
        raise TemplateError(
            "No undecorated layout found in the template — every layout carries "
            "background art that would render behind slide content. Add a truly "
            f"blank layout to brand/template.pptx. (Least decorated: '{layout.name}' "
            f"with {decorations} graphic(s).)")
    return layout


def _is_graphic(shape) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)


def add_blank_slide(prs):
    """Add a slide on the blank layout and remove any inherited placeholders."""
    slide = prs.slides.add_slide(blank_layout(prs))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide
