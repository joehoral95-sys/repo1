"""Rich text helpers: light inline markup -> styled runs.

Spec strings may contain:
  **bold**              -> bold run
  [accent]text[/accent] -> brand accent color
These are the ONLY styling hooks a spec gets; everything else is engine-owned.
"""

from __future__ import annotations

import re
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from ..tokens import Tokens

_TOKEN_RE = re.compile(r"(\*\*.+?\*\*|\[accent\].+?\[/accent\])", re.DOTALL)


def adaptive_pt(base_pt: int, texts: str | list[str],
                tiers: list[tuple[int, int]]) -> int:
    """Scale type UP when content is short (never below base_pt).

    tiers: (max_chars, pt) pairs, ascending by max_chars - the first tier
    whose max_chars covers the LONGEST text wins. Example:
    adaptive_pt(14, labels, [(14, 20), (22, 17)]) → 20pt when every label is
    ≤14 chars, 17pt when ≤22, else 14pt. Markup is stripped before counting.
    """
    if isinstance(texts, str):
        texts = [texts]
    longest = max((len(strip_markup(t)) for t in texts), default=0)
    for max_chars, pt in tiers:
        if longest <= max_chars:
            return max(pt, base_pt)
    return base_pt


@dataclass
class Segment:
    text: str
    bold: bool = False
    accent: bool = False


def parse_markup(text: str) -> list[Segment]:
    segments: list[Segment] = []
    for part in _TOKEN_RE.split(text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            segments.append(Segment(part[2:-2], bold=True))
        elif part.startswith("[accent]") and part.endswith("[/accent]"):
            segments.append(Segment(part[len("[accent]"):-len("[/accent]")], accent=True))
        else:
            segments.append(Segment(part))
    return segments


def strip_markup(text: str) -> str:
    return "".join(seg.text for seg in parse_markup(text))


def set_runs(paragraph, text: str, *, font: str, size: Pt, color: RGBColor,
             accent_color: RGBColor, bold: bool = False,
             letter_spacing_pt: float | None = None) -> None:
    """Fill a paragraph with runs parsed from markup, clearing existing runs."""
    # paragraph may come with an empty default run; clear text first
    for run in list(paragraph.runs):
        run._r.getparent().remove(run._r)
    for seg in parse_markup(text):
        run = paragraph.add_run()
        run.text = seg.text
        run.font.name = font
        run.font.size = size
        run.font.bold = bold or seg.bold
        run.font.color.rgb = accent_color if seg.accent else color
        if letter_spacing_pt:
            # DrawingML spc is in 1/100 pt
            run.font._rPr.set("spc", str(int(letter_spacing_pt * 100)))


def add_text(
    slide,
    box,
    text: str,
    tokens: Tokens,
    *,
    scale: str = "body",
    role: str = "body",          # 'heading' or 'body' -> font family
    color: str = "neutral_dark",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
    shrink_to_fit: bool = False,
    letter_spacing_pt: float | None = None,
    size_pt: int | None = None,
):
    """Add a text box at `box` with brand styling. Returns the shape.

    size_pt overrides the scale's point size (used by adaptive sizing)."""
    shape = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    font = tokens.heading_font.name if role == "heading" else tokens.body_font.name
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        set_runs(p, line, font=font,
                 size=Pt(size_pt) if size_pt else tokens.pt(scale),
                 color=tokens.color(color), accent_color=tokens.color("accent"),
                 bold=bold, letter_spacing_pt=letter_spacing_pt)
    if shrink_to_fit:
        _enable_shrink(tf)
    return shape


def add_bullets(
    slide,
    box,
    bullets,               # list[str | (str, list[str])]
    tokens: Tokens,
    *,
    scale: str = "body",
    color: str = "neutral_dark",
    space_after_pt: int = 10,
    size_pt: int | None = None,
):
    """Bulleted text block. Accepts strings or (text, sub_items) tuples.

    size_pt overrides the top-level bullet size (adaptive sizing)."""
    shape = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    font = tokens.body_font.name
    first = True

    def _para(text: str, level: int, size_key: str):
        nonlocal first
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(space_after_pt if level == 0 else max(4, space_after_pt - 4))
        p.line_spacing = 1.1
        size = Pt(size_pt) if (size_pt and level == 0) else tokens.pt(size_key)
        set_runs(p, text, font=font, size=size,
                 color=tokens.color(color), accent_color=tokens.color("accent"))
        _apply_bullet_char(p, level, tokens)

    for item in bullets:
        if isinstance(item, tuple):
            text, subs = item
            _para(text, 0, scale)
            for s in subs:
                _para(s, 1, "bullet_sub")
        else:
            _para(item, 0, scale)
    return shape


def _apply_bullet_char(paragraph, level: int, tokens: Tokens) -> None:
    """Give paragraphs a clean brand bullet (en dash level 0, · level 1)."""
    from lxml import etree

    pPr = paragraph._p.get_or_add_pPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ("buNone", "buChar", "buAutoNum"):
        for el in pPr.findall(f"{ns}{tag}"):
            pPr.remove(el)
    buFont = etree.SubElement(pPr, f"{ns}buFont")
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, f"{ns}buChar")
    buChar.set("char", "–" if level == 0 else "·")
    pPr.set("indent", "-182880")  # 0.2in hanging indent
    pPr.set("marL", str(182880 if level == 0 else 457200))


def _enable_shrink(text_frame) -> None:
    """Turn on PowerPoint's shrink-text-on-overflow for a text frame."""
    from lxml import etree

    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    bodyPr = text_frame._txBody.find(f"{ns}bodyPr")
    for tag in ("spAutoFit", "noAutofit", "normAutofit"):
        for el in bodyPr.findall(f"{ns}{tag}"):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, f"{ns}normAutofit")
