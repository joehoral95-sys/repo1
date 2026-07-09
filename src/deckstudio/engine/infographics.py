"""Composed-shape infographic components.

Used when data is 1-4 numbers and a default chart would underserve it.
Everything here is native shapes and text — editable in PowerPoint like any
other object. (These are not data-linked like charts; the spec is the edit
path for regenerating them.)
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from ..spec.schema import Stat
from ..tokens import Tokens
from .geometry import Box, inset
from .shapes import add_rect
from .text import add_text

ARROWS = {"up": "\u25b2", "down": "\u25bc", "flat": "\u25ac", "none": ""}
SENTIMENT_COLOR = {"good": "positive", "bad": "negative", "neutral": "neutral_mid"}
STATUS_COLOR = {"on_track": "positive", "watch": "warning", "pending": "neutral_mid"}


def add_stat_tile(slide, box: Box, stat: Stat, tokens: Tokens, *, hero: bool = False,
                  accent_index: int = 0) -> None:
    """A KPI status card (the SOA board-deck pattern): metric label on top,
    big value, delta and target context lines, status dot in the corner.

    hero=True renders the card on the primary brand color (the headline stat).
    """
    del accent_index  # cards differentiate via status dots now
    if hero:
        add_rect(slide, box, tokens, fill="primary", rounded=True, corner=0.06)
    else:
        add_rect(slide, box, tokens, fill="neutral_light", rounded=True, corner=0.06,
                 outline="border")
    label_color = "accent_warm" if hero else "neutral_mid"
    value_color = "white" if hero else "primary"

    if stat.status:
        dot_d = 0.14
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(box.right_in - 0.32 - dot_d),
            Inches(box.top_in + 0.26), Inches(dot_d), Inches(dot_d))
        dot.fill.solid()
        dot.fill.fore_color.rgb = tokens.color(STATUS_COLOR[stat.status])
        dot.line.fill.background()
        dot.shadow.inherit = False

    pad = inset(box, x_in=0.3, y_in=0.24)
    add_text(slide, Box(pad.left_in, pad.top_in, pad.width_in - 0.25, 0.55),
             stat.label.upper(), tokens, scale="stat_label",
             color=label_color, bold=True, letter_spacing_pt=1.0,
             shrink_to_fit=True)

    value_top = pad.top_in + box.height_in * 0.24
    value_h = box.height_in * 0.34
    scale = "big_number" if box.width_in > 2.6 else "section_title"
    add_text(slide, Box(pad.left_in, value_top, pad.width_in, value_h),
             stat.value, tokens, scale=scale, role="heading",
             color=value_color, bold=True, anchor=MSO_ANCHOR.MIDDLE,
             shrink_to_fit=True)

    line_y = value_top + value_h + 0.14
    if stat.delta:
        delta_text = f"{ARROWS[stat.arrow]} {stat.delta}".strip()
        delta_color = SENTIMENT_COLOR[stat.sentiment]
        if hero:
            delta_color = "accent_warm" if stat.sentiment != "neutral" else "white"
        add_text(slide, Box(pad.left_in, line_y, pad.width_in, 0.26), delta_text,
                 tokens, scale="stat_delta", color=delta_color, bold=True)
        line_y += 0.32
    if stat.target:
        add_text(slide, Box(pad.left_in, line_y, pad.width_in, 0.26), stat.target,
                 tokens, scale="caption",
                 color="neutral_light" if hero else "neutral_mid",
                 shrink_to_fit=True)


def add_progress_ring(slide, box: Box, fraction: float, label: str, tokens: Tokens) -> None:
    """A donut-style progress indicator built from native shapes.

    Track = full circle in neutral_light; progress = PIE shape in accent;
    center disc carries the percentage text.
    """
    fraction = max(0.0, min(1.0, fraction))
    size = min(box.width_in, box.height_in)
    cx = box.left_in + (box.width_in - size) / 2
    cy = box.top_in + (box.height_in - size) / 2
    circle_box = Box(cx, cy, size, size)

    track = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, circle_box.left, circle_box.top, circle_box.width, circle_box.height)
    track.fill.solid()
    track.fill.fore_color.rgb = tokens.color("neutral_light")
    track.line.fill.background()
    track.shadow.inherit = False

    if fraction > 0:
        pie = slide.shapes.add_shape(
            MSO_SHAPE.PIE, circle_box.left, circle_box.top, circle_box.width, circle_box.height)
        # PIE adjustments are start/end angles in degrees; 0 = 3 o'clock.
        # Start at 12 o'clock (-90) and sweep clockwise.
        try:
            pie.adjustments[0] = -90.0
            pie.adjustments[1] = -90.0 + 360.0 * fraction
        except (IndexError, ValueError):
            pass
        pie.fill.solid()
        pie.fill.fore_color.rgb = tokens.color("accent")
        pie.line.fill.background()
        pie.shadow.inherit = False

    hole_scale = 0.62
    hole_size = size * hole_scale
    hole_off = (size - hole_size) / 2
    hole = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx + hole_off), Inches(cy + hole_off),
        Inches(hole_size), Inches(hole_size))
    hole.fill.solid()
    hole.fill.fore_color.rgb = tokens.color("white")
    hole.line.fill.background()
    hole.shadow.inherit = False

    pct_box = Box(cx + hole_off, cy + hole_off + hole_size * 0.18, hole_size, hole_size * 0.4)
    add_text(slide, pct_box, f"{round(fraction * 100)}%", tokens, scale="section_title",
             role="heading", color="primary", bold=True, align=PP_ALIGN.CENTER,
             shrink_to_fit=True)
    label_box = Box(cx + hole_off, pct_box.bottom_in, hole_size, hole_size * 0.28)
    add_text(slide, label_box, label, tokens, scale="stat_label", color="neutral_mid",
             align=PP_ALIGN.CENTER, shrink_to_fit=True)


def add_proportion_bars(slide, box: Box, items: list[tuple[str, float]],
                        tokens: Tokens) -> None:
    """Horizontal proportional bars: label + bar scaled to the max value."""
    if not items:
        return
    max_val = max(v for _, v in items) or 1.0
    row_h = box.height_in / len(items)
    bar_h = min(0.32, row_h * 0.45)
    label_w = box.width_in * 0.28
    bar_area_w = box.width_in - label_w - 0.15
    for i, (label, value) in enumerate(items):
        top = box.top_in + i * row_h + (row_h - bar_h) / 2
        add_text(slide, Box(box.left_in, top - 0.02, label_w, bar_h + 0.1), label,
                 tokens, scale="stat_label", color="neutral_dark",
                 anchor=MSO_ANCHOR.MIDDLE, shrink_to_fit=True)
        width = max(0.05, bar_area_w * (value / max_val))
        color = "accent" if value == max_val else "primary"
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box.left_in + label_w + 0.15), Inches(top),
            Inches(width), Inches(bar_h))
        bar.adjustments[0] = 0.5
        bar.fill.solid()
        bar.fill.fore_color.rgb = tokens.color(color)
        bar.line.fill.background()
        bar.shadow.inherit = False
