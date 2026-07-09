"""Shape helpers shared by renderers: brand-styled rectangles, accent rules, images."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from ..tokens import Tokens
from .geometry import Box


def add_rect(slide, box: Box, tokens: Tokens, *, fill: str | None = "neutral_light",
             rounded: bool = True, corner: float = 0.08, outline: str | None = None,
             shadow: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, box.left, box.top, box.width, box.height)
    if rounded:
        try:
            shape.adjustments[0] = corner
        except (IndexError, ValueError):
            pass
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = tokens.color(fill)
    else:
        shape.fill.background()
    if outline:
        shape.line.color.rgb = tokens.color(outline)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left_in: float, top_in: float, width_in: float,
                   tokens: Tokens, *, color: str = "accent", height_in: float = 0.06):
    """The thin brand rule used under titles and on section slides."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = tokens.color(color)
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def fill_background(slide, tokens: Tokens, color: str) -> None:
    """Dark slide background with a subtle diagonal gradient (deep navy
    lifting toward blue) — reads richer than a flat fill. Falls back to
    solid if the gradient API is unavailable."""
    bg = slide.background
    if color == "primary":
        try:
            bg.fill.gradient()
            stops = bg.fill.gradient_stops
            stops[0].color.rgb = RGBColor.from_string("0A1F33")
            stops[0].position = 0.0
            stops[1].color.rgb = RGBColor.from_string("15355C")
            stops[1].position = 1.0
            bg.fill.gradient_angle = 65.0
            return
        except Exception:
            pass
    bg.fill.solid()
    bg.fill.fore_color.rgb = tokens.color(color)


def add_logo(slide, tokens: Tokens, *, dark_bg: bool = False,
             left_in: float | None = None, center_at_in: float | None = None,
             bottom_in: float = 7.18, height_in: float = 0.32):
    """Place the brand logo BOTTOM-LEFT (the SOA template convention), or
    centered on `center_at_in` for symmetric compositions; silently skip if
    no asset exists."""
    path = tokens.logo_dark_bg if dark_bg else tokens.logo_default
    if not path:
        return None
    left = Inches(left_in if left_in is not None else tokens.margin_in)
    pic = slide.shapes.add_picture(str(path), left, Emu(0), height=Inches(height_in))
    if center_at_in is not None:
        pic.left = Inches(center_at_in) - pic.width // 2
    pic.top = Inches(bottom_in) - pic.height
    return pic


def chip_width(text: str) -> float:
    """Estimated pill width (inches) for layout math at chip point size."""
    return 0.125 * len(text) + 0.38


def add_chip(slide, left_in: float, top_in: float, text: str, tokens: Tokens, *,
             fill: str = "accent_warm", text_color: str = "primary",
             height_in: float = 0.36):
    """A small pill badge ('RECOMMENDED', '▼ 7.1% vs Q1'). Returns the shape."""
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    from .text import add_text

    width_in = chip_width(text)
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = tokens.color(fill)
    pill.line.fill.background()
    pill.shadow.inherit = False
    add_text(slide, Box(left_in, top_in + 0.02, width_in, height_in - 0.04),
             text, tokens, scale="stat_delta", color=text_color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return pill


def add_brand_art(slide, tokens: Tokens, name: str, box: Box, *, stretch: bool = False):
    """Place a brand accent graphic (brand/assets/patterns/<name>.png) inside
    `box`. stretch=True fills the box exactly (how the SOA template applies
    its hand-drawn marks); otherwise fit preserving aspect ratio. Silently
    skips if the asset doesn't exist — art is garnish, never load-bearing."""
    path = tokens.brand_dir / "assets" / "patterns" / f"{name}.png"
    if not path.exists():
        return None
    if stretch:
        return slide.shapes.add_picture(str(path), box.left, box.top,
                                        width=box.width, height=box.height)
    return add_picture_fitted(slide, str(path), box)


def add_picture_fitted(slide, path: str, box: Box):
    """Place a picture centered inside `box`, preserving aspect ratio."""
    from PIL import Image

    with Image.open(path) as im:
        w, h = im.size
    box_ratio = box.width_in / box.height_in
    img_ratio = w / h
    if img_ratio > box_ratio:
        width = box.width
        pic = slide.shapes.add_picture(path, box.left, box.top, width=width)
        pic.top = box.top + Emu(int((box.height - pic.height) / 2))
    else:
        height = box.height
        pic = slide.shapes.add_picture(path, box.left, box.top, height=height)
        pic.left = box.left + Emu(int((box.width - pic.width) / 2))
    return pic
