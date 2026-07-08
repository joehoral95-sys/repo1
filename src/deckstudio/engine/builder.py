"""The spec -> pptx compiler. Deterministic: same spec + tokens + template
in, same deck out. All styling decisions live here and in renderers/tokens —
never in the spec.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from ..spec.schema import DeckSpec
from ..tokens import Tokens, load_tokens
from . import animations
from .geometry import SLIDE_H_IN, SLIDE_W_IN, Box
from .registry import get_renderer
from .template import add_blank_slide, open_template
from .text import add_text


@dataclass
class RenderContext:
    prs: object
    tokens: Tokens
    spec: DeckSpec
    slide_index: int = 0          # 0-based index of the slide being rendered
    warnings: list[str] = field(default_factory=list)

    @property
    def deck(self):
        return self.spec.deck

    def warn(self, message: str) -> None:
        self.warnings.append(message)


@dataclass
class BuildResult:
    output_path: Path
    version: int
    slide_count: int
    warnings: list[str]


# Slide types that own their whole canvas: no footer/page number drawn on them.
FULL_BLEED_TYPES = {"title", "section", "quote", "thanks"}


def build_deck(
    spec: DeckSpec,
    *,
    tokens: Tokens | None = None,
    template_path: Path | None = None,
    output_dir: Path,
    deck_name: str,
    enable_animations: bool = True,
) -> BuildResult:
    tokens = tokens or load_tokens()
    if template_path is None:
        default_template = tokens.brand_dir / "template.pptx"
        template_path = default_template if default_template.exists() else None
    prs = open_template(template_path)

    ctx = RenderContext(prs=prs, tokens=tokens, spec=spec)
    animate_deck = enable_animations and spec.deck.animations != "off"

    for i, slide_model in enumerate(spec.slides):
        ctx.slide_index = i
        slide = add_blank_slide(prs)
        render = get_renderer(slide_model.type)
        render(slide, slide_model, ctx)

        if slide_model.type not in FULL_BLEED_TYPES:
            _add_footer(slide, ctx, page_number=i + 1)

        if slide_model.notes:
            slide.notes_slide.notes_text_frame.text = slide_model.notes

        if animate_deck and slide_model.animate != "none":
            try:
                animations.apply(slide, mode=slide_model.animate)
            except Exception as e:  # animations must never break a build
                ctx.warn(f"slide '{slide_model.id}': animation skipped ({e})")

    output_dir.mkdir(parents=True, exist_ok=True)
    version = _next_version(output_dir, deck_name)
    out_path = output_dir / f"{deck_name}_v{version}.pptx"
    prs.save(str(out_path))
    return BuildResult(
        output_path=out_path,
        version=version,
        slide_count=len(spec.slides),
        warnings=ctx.warnings,
    )


def _add_footer(slide, ctx: RenderContext, page_number: int) -> None:
    tokens = ctx.tokens
    text = ctx.deck.footer if ctx.deck.footer is not None else tokens.footer_text
    footer_box = Box(tokens.margin_in, SLIDE_H_IN - 0.42,
                     SLIDE_W_IN - 2 * tokens.margin_in - 0.5, 0.3)
    if text:
        add_text(slide, footer_box, text, tokens, scale="footer", color="neutral_mid")
    if tokens.show_page_numbers:
        num_box = Box(SLIDE_W_IN - tokens.margin_in - 0.5, SLIDE_H_IN - 0.42, 0.5, 0.3)
        from pptx.enum.text import PP_ALIGN
        add_text(slide, num_box, str(page_number), tokens, scale="footer",
                 color="neutral_mid", align=PP_ALIGN.RIGHT)


_VERSION_RE = re.compile(r"_v(\d+)\.pptx$")


def _next_version(output_dir: Path, deck_name: str) -> int:
    """Outputs are versioned and never overwritten: deck_v1.pptx, deck_v2.pptx, ..."""
    latest = 0
    for p in output_dir.glob(f"{deck_name}_v*.pptx"):
        m = _VERSION_RE.search(p.name)
        if m and p.name == f"{deck_name}_v{m.group(1)}.pptx":
            latest = max(latest, int(m.group(1)))
    return latest + 1
