"""Native PowerPoint charts, restyled to brand.

THE RULE: charts are created exclusively with python-pptx's add_chart(), which
embeds a real editable workbook - Joe can right-click -> Edit Data in
PowerPoint and the graph updates. Never render a chart as an image.
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.util import Pt

from ..spec.schema import ChartSpec
from ..tokens import Tokens
from .geometry import Box

_KIND_TO_XL = {
    "bar_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

_LIGHT_GRID = RGBColor.from_string("E3E8ED")


def add_chart(slide, spec: ChartSpec, box: Box, tokens: Tokens):
    """Create and brand-style a native chart. Returns the GraphicFrame."""
    xl_type = _KIND_TO_XL[spec.kind]

    if spec.kind == "scatter":
        data = XyChartData()
        for s in spec.series:
            series_data = data.add_series(s.name)
            for cat, val in zip(spec.categories, s.values, strict=True):
                if val is not None:
                    series_data.add_data_point(float(cat), float(val))
    else:
        data = CategoryChartData()
        data.categories = spec.categories
        for s in spec.series:
            data.add_series(s.name, s.values)

    frame = slide.shapes.add_chart(xl_type, box.left, box.top, box.width, box.height, data)
    chart = frame.chart
    _style_chart(chart, spec, tokens)
    return frame


def _style_chart(chart, spec: ChartSpec, tokens: Tokens) -> None:
    body_font = tokens.body_font.name
    chart.font.name = body_font
    chart.font.size = Pt(tokens.chart_font_size)
    chart.font.color.rgb = tokens.color("neutral_dark")
    chart.has_title = False  # the slide title carries the assertion

    # Legend only when there is something to disambiguate.
    multi = len(spec.series) > 1 or spec.kind in ("pie", "doughnut")
    chart.has_legend = multi
    if multi:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(tokens.chart_font_size)

    if spec.kind in ("pie", "doughnut"):
        _color_slices(chart, spec, tokens)
    else:
        _color_series(chart, spec, tokens)
        _style_axes(chart, spec, tokens)

    _apply_data_labels(chart, spec, tokens)
    if spec.highlight is not None:
        _apply_highlight(chart, spec, tokens)


def _color_series(chart, spec: ChartSpec, tokens: Tokens) -> None:
    palette = tokens.chart_series
    line_kind = spec.kind == "line"
    for i, series in enumerate(chart.series):
        color = palette[i % len(palette)]
        if line_kind:
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2.5)
            try:
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = color
                series.marker.format.line.fill.background()
            except (AttributeError, ValueError):
                pass
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.fill.background()


def _color_slices(chart, spec: ChartSpec, tokens: Tokens) -> None:
    palette = tokens.chart_series
    series = chart.plots[0].series[0]
    for i in range(len(spec.categories)):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = palette[i % len(palette)]


def _style_axes(chart, spec: ChartSpec, tokens: Tokens) -> None:
    grid_color = _LIGHT_GRID
    try:
        vaxis = chart.value_axis
        if spec.kind in ("bar_clustered", "bar_stacked", "bar_horizontal", "area"):
            # honesty rule: bars/areas always start at zero - no dramatic
            # differences faked by a truncated axis
            vaxis.minimum_scale = 0.0
        vaxis.has_major_gridlines = True
        vaxis.major_gridlines.format.line.color.rgb = grid_color
        vaxis.major_gridlines.format.line.width = Pt(0.5)
        vaxis.has_minor_gridlines = False
        vaxis.format.line.fill.background()
        vaxis.major_tick_mark = XL_TICK_MARK.NONE
        vaxis.tick_labels.font.size = Pt(tokens.chart_font_size - 1)
        vaxis.tick_labels.font.color.rgb = tokens.color("neutral_mid")
        if spec.number_format:
            vaxis.tick_labels.number_format = spec.number_format
            vaxis.tick_labels.number_format_is_linked = False
    except ValueError:
        pass
    try:
        caxis = chart.category_axis
        caxis.has_major_gridlines = False
        caxis.format.line.color.rgb = grid_color
        caxis.major_tick_mark = XL_TICK_MARK.NONE
        caxis.tick_labels.font.size = Pt(tokens.chart_font_size)
        caxis.tick_labels.font.color.rgb = tokens.color("neutral_dark")
    except ValueError:
        pass


def _apply_data_labels(chart, spec: ChartSpec, tokens: Tokens) -> None:
    """Data labels where they help: single-series bars and pie/doughnut."""
    from pptx.enum.chart import XL_LABEL_POSITION

    plot = chart.plots[0]
    if spec.kind in ("pie", "doughnut"):
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.show_percentage = True
        labels.show_value = False
        labels.number_format = "0%"
        labels.number_format_is_linked = False
        labels.font.size = Pt(tokens.chart_font_size)
        labels.font.color.rgb = tokens.color("white")
        if spec.kind == "pie":
            labels.position = XL_LABEL_POSITION.INSIDE_END
    elif spec.kind in ("bar_clustered", "bar_horizontal") and len(spec.series) == 1:
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.position = XL_LABEL_POSITION.OUTSIDE_END
        labels.font.size = Pt(tokens.chart_font_size - 1)
        labels.font.color.rgb = tokens.color("neutral_dark")
        if spec.number_format:
            labels.number_format = spec.number_format
            labels.number_format_is_linked = False


def _apply_highlight(chart, spec: ChartSpec, tokens: Tokens) -> None:
    """Recolor one data point to SOA yellow - the 'look here' cue (distinct
    from every series color in the palette)."""
    accent = tokens.color("highlight")
    idx = next(i for i, s in enumerate(spec.series) if s.name == spec.highlight.series)
    series = chart.plots[0].series[idx] if spec.kind in ("pie", "doughnut") else chart.series[idx]
    point = series.points[spec.highlight.point]
    if spec.kind == "line":
        point.marker.format.fill.solid()
        point.marker.format.fill.fore_color.rgb = accent
    else:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = accent
