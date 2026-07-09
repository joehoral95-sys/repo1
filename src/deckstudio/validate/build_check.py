"""Post-build smoke check: reopen the produced deck and assert invariants.

Runs automatically after every build. Failures here mean an engine bug, not a
spec problem - they should be reported, not 'fixed' by editing the spec.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..spec.schema import DeckSpec


def check_build(pptx_path: Path, spec: DeckSpec) -> list[str]:
    """Return a list of problems (empty = healthy)."""
    problems: list[str] = []
    prs = Presentation(str(pptx_path))

    if len(prs.slides) != len(spec.slides):
        problems.append(
            f"slide count mismatch: spec has {len(spec.slides)}, "
            f"output has {len(prs.slides)}")

    expected_charts = sum(1 for s in spec.slides if s.type == "chart")
    found_charts = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART or getattr(shape, "has_chart", False):
                found_charts += 1
                chart_part = shape.chart.part
                # THE editability guarantee: every chart must carry an
                # embedded workbook so Edit Data works in PowerPoint.
                xlsx_part = getattr(chart_part.chart_workbook, "xlsx_part", None)
                if xlsx_part is None:
                    problems.append(
                        f"chart on slide {prs.slides.index(slide) + 1} has no embedded "
                        "workbook - Edit Data would not work")
    if found_charts != expected_charts:
        problems.append(
            f"chart count mismatch: spec has {expected_charts}, output has {found_charts}")

    size = pptx_path.stat().st_size
    if size < 10_000:
        problems.append(f"output file is suspiciously small ({size} bytes)")

    return problems
